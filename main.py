import requests, os
import csv
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from selenium import webdriver
from selenium.webdriver.chrome.service import Service
from selenium.webdriver.common.by import By
from selenium.webdriver.support.ui import WebDriverWait
from selenium.webdriver.support import expected_conditions as EC
from selenium.webdriver.chrome.options import Options
from bs4 import BeautifulSoup
from tqdm import tqdm


api_key_vagalume = os.getenv("KEY_VAGALUME")

def get_lyrics():

    lyrics_list = []
    with open('musicas.csv', 'r') as input_musics:
        lines = input_musics.readlines()
        for line in tqdm(lines, total=len(lines), desc="Processando as letras"):
            music, artist_uf = line.split(';')
            artist = artist_uf.replace('\n', '')
            local_lyric = get_lyric_local(music, artist)
            if local_lyric:
                lyrics_list.append(local_lyric)
            else:
                api_lyric = get_lyrics_on_vagalume_api(music, artist)
                if api_lyric:
                    lyrics_list.append(api_lyric)
                else:
                    scrapping_lyric = get_lyric_by_scraping(music, artist)
                    if scrapping_lyric:
                        lyrics_list.append(scrapping_lyric)
                    else:
                        lyrics_list.append(f'Música {music} do artista {artist} não encontrada')
    return lyrics_list

def get_lyric_local(music, artist):
    with open('lista_com_id.csv', 'r') as local_list:
        for line in local_list.readlines():
            music_name, artist_name, lyric = line.split(';')
            if music == music_name and artist == artist_name:
                return lyric

def get_lyrics_on_vagalume_api(music, artist):
    url = f"https://api.vagalume.com.br/search.php?art={artist}&mus={music}&apikey={api_key_vagalume}"
    response = requests.get(url)
    if response.status_code == 200:
        data = response.json()
        # print(data)
        if 'mus' in data:
            lyric = data['mus'][0]['text']
            formated_lyric = lyric.replace('\n', ' | ')
            with open("lista_com_id.csv", "a", newline='') as local_list:
                writer = csv.writer(local_list, delimiter=';')
                row = [music, artist, formated_lyric]
                writer.writerow(row)
            return formated_lyric
    return None

def get_lyric_by_scraping(music_name, artist):
    chrome_options = Options()
    chrome_options.add_argument("--headless")  # Executa o Chrome em modo invisível (sem interface gráfica)
    service = Service(executable_path='/usr/local/bin/chromedriver')
    # Use with block para gerenciar o driver
    with webdriver.Chrome(service=service, options=chrome_options) as driver:
        # Define a query de pesquisa
        search_query = f"{artist} {music_name}"
        url = f"https://www.letras.mus.br/?q={search_query.replace(' ', '%20')}"
        # Navega até a URL
        driver.get(url)
        try:
            # Aguarde até que a classe `gs-title` seja carregada
            WebDriverWait(driver, 10).until(EC.presence_of_element_located((By.CSS_SELECTOR, 'div.gs-title > a')))
            # Encontre o primeiro elemento `a` dentro da `div.gs-title`
            result = driver.find_element(By.CSS_SELECTOR, 'div.gs-title > a')
            link = result.get_attribute('href')  # Pega o link da música
            # Navega até a página da música
            driver.get(link)
            # Aguarde o carregamento da letra
            WebDriverWait(driver, 10).until(EC.presence_of_element_located((By.CLASS_NAME, 'lyric-original')))
            result = driver.find_element(By.CLASS_NAME, 'lyric-original')
            # Use BeautifulSoup para processar o HTML e extrair a letra
            soup = BeautifulSoup(result.get_attribute('outerHTML'), 'html.parser')
            lyric = soup.get_text(separator=" | ").strip()
            if lyric:
                with open("lista_com_id.csv", "a", newline='') as local_list:
                    writer = csv.writer(local_list, delimiter=';')
                    row = [music_name, artist, lyric]
                    writer.writerow(row)
                return lyric
            else:
                with open("lista_com_id.csv", "a") as id_list:
                    id_list.write(f"{music_name};{artist};Letra não encontrada no letras.mus\n")
                    return None
        except Exception as e:
            return f"Error:{str(e)}"

# Função para dividir o texto em múltiplos slides de até 300 caracteres por slide
def split_lyric_into_slides(lyrics, max_chars=300):
    slides = []
    current_slide = ""
    words = lyrics.split()

    for word in words:
        # Adiciona o próximo trecho ao slide atual
        if len(current_slide) + len(word) + 1 <= max_chars:
            current_slide += " " + word
        else:
            # Insere os " | " de volta antes de adicionar o slide à lista
            slides.append(current_slide.strip() + " | ")
            current_slide = word

    # Adiciona o último slide se houver conteúdo restante
    if current_slide:
        slides.append(current_slide.strip() + " | ")

    return slides

# Função para criar a apresentação PPTX
def create_presentation_with_format(lyrics_list):
    # Cria uma nova apresentação
    prs = Presentation()

    # Configura o tamanho do slide para 25,4 cm x 19,05 cm (4:3)
    prs.slide_width = Inches(10)  # 25,4 cm = 10 inches
    prs.slide_height = Inches(7.5)  # 19,05 cm = 7.5 inches

    for lyric in tqdm(lyrics_list, total=len(lyrics_list), desc="Montando apresentação"):
        slides = split_lyric_into_slides(lyric)
        for part in slides:
            # Adiciona um slide branco (layout [6])
            slide_layout = prs.slide_layouts[6]  # Slide branco
            slide = prs.slides.add_slide(slide_layout)

            # Adiciona uma caixa de texto centralizada
            left = top = Inches(0)  # Margens 0 cm
            width = prs.slide_width
            height = prs.slide_height

            textbox = slide.shapes.add_textbox(left, top, width, height)
            text_frame = textbox.text_frame
            text_frame.word_wrap = True
            text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE  # Centralizado verticalmente

            # Divide a letra onde há " | " e adiciona cada parte como um parágrafo separado
            lyrics_lines = [line.strip() for line in part.split("|") if line.strip()]

            if lyrics_lines:
                # Define o primeiro parágrafo diretamente no `text_frame.text` para evitar a linha vazia
                text_frame.text = lyrics_lines[0].upper()
                p = text_frame.paragraphs[0]
                p.font.size = Pt(36)  # Fonte tamanho 36pt
                p.font.name = 'Arial'
                p.alignment = PP_ALIGN.CENTER

                # Adiciona os parágrafos seguintes
                for line in lyrics_lines[1:]:
                    p = text_frame.add_paragraph()
                    p.text = line.upper()  # Coloca o texto em maiúsculas
                    p.font.size = Pt(36)  # Fonte tamanho 36pt
                    p.font.name = 'Arial'
                    p.alignment = PP_ALIGN.CENTER  # Centralizado horizontalmente

    # Salva a apresentação
    prs.save('apresentacao_letras.pptx')

if __name__ == '__main__':
    lyrics_list = get_lyrics()
    create_presentation_with_format(lyrics_list)
