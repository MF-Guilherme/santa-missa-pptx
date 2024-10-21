from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR  # Import correto do MSO_ANCHOR
from pptx.dml.color import RGBColor


# Função para dividir o texto em múltiplos slides de até 300 caracteres por slide
def split_lyric_into_slides(lyrics, max_chars=300):
    slides = []
    current_slide = ""
    words = lyrics.split()

    for word in words:
        # Adiciona o próximo trecho ao slide atual
        if len(current_slide) + len(word) + 1 <= max_chars:
            current_slide += " " + word
        else:
            # Insere os " | " de volta antes de adicionar o slide à lista
            slides.append(current_slide.strip() + " | ")
            current_slide = word

    # Adiciona o último slide se houver conteúdo restante
    if current_slide:
        slides.append(current_slide.strip() + " | ")

    return slides



# Função para criar a apresentação PPTX
def create_presentation_with_format(lyrics_list):
    # Cria uma nova apresentação
    prs = Presentation()

    # Configura o tamanho do slide para 25,4 cm x 19,05 cm (4:3)
    prs.slide_width = Inches(10)  # 25,4 cm = 10 inches
    prs.slide_height = Inches(7.5)  # 19,05 cm = 7.5 inches

    for lyric in lyrics_list:
        slides = split_lyric_into_slides(lyric)
        for part in slides:
            # Adiciona um slide branco (layout [6])
            slide_layout = prs.slide_layouts[6]  # Slide branco
            slide = prs.slides.add_slide(slide_layout)

            # Adiciona uma caixa de texto centralizada
            left = top = Inches(0)  # Margens 0 cm
            width = prs.slide_width
            height = prs.slide_height

            textbox = slide.shapes.add_textbox(left, top, width, height)
            text_frame = textbox.text_frame
            text_frame.word_wrap = True
            text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE  # Centralizado verticalmente

            # Divide a letra onde há " | " e adiciona cada parte como um parágrafo separado
            lyrics_lines = [line.strip() for line in part.split("|") if line.strip()]

            if lyrics_lines:
                # Define o primeiro parágrafo diretamente no `text_frame.text` para evitar a linha vazia
                text_frame.text = lyrics_lines[0].upper()
                p = text_frame.paragraphs[0]
                p.font.size = Pt(36)  # Fonte tamanho 36pt
                p.font.name = 'Arial'
                p.alignment = PP_ALIGN.CENTER

                # Adiciona os parágrafos seguintes
                for line in lyrics_lines[1:]:
                    p = text_frame.add_paragraph()
                    p.text = line.upper()  # Coloca o texto em maiúsculas
                    p.font.size = Pt(36)  # Fonte tamanho 36pt
                    p.font.name = 'Arial'
                    p.alignment = PP_ALIGN.CENTER  # Centralizado horizontalmente

    # Salva a apresentação
    prs.save('apresentacao_letras.pptx')

# Exemplo de uso
lyric = ["Uma vida de graça e vitória | É o que Deus tem pra te dar | A condição pra você recebê-la | É a busca da Santidade |  | Uma vida de graça e vitória | É o que Deus tem pra te dar | A condição pra você recebê-la | É a busca da Santidade |  | Lutar contra todo pecado | Que te separa de Deus | Romper com todas cadeias | Ser livre,inteiramente de Deus |  | Eu quero ser de Deus | Eu quero ser de Deus | Eu quero viver esse amor! |  | Eu quero ser de Deus | Eu quero ser de Deus | Eu quero viver esse amor |  | Ooh,ooh,ooh! | Ooh,ooh,ooh!."]

create_presentation_with_format(lyric)
